"""
Build graduation defense PPT v2 using 进度.pptx template theme/master.
Creates a NEW file: 毕业答辩-聂俊宇-v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import os

TEMPLATE = r"D:\GraduationProject\Doc\PPT\进度.pptx"
OUTPUT = r"D:\GraduationProject\Doc\PPT\毕业答辩-聂俊宇-v2.pptx"

prs = Presentation(TEMPLATE)

# Map layout names to indices
layouts = {}
for i, lay in enumerate(prs.slide_layouts):
    layouts[lay.name] = i
    print(f"Layout {i}: {lay.name}")

print(f"\nSlide size: {prs.slide_width/914400:.1f}x{prs.slide_height/914400:.1f} in")
print(f"Existing slides: {len(prs.slides)}")

# We'll KEEP the first 4 slides as reference and append new ones
# OR start fresh with same template master

# Strategy: create a fresh presentation using the same template
prs2 = Presentation(TEMPLATE)
# Delete all existing slides (we'll create our own)
# We need to work with the slide master from this template
master = prs2.slide_masters[0]

# Re-map layouts in prs2
layouts2 = {}
for i, lay in enumerate(prs2.slide_layouts):
    layouts2[lay.name] = i

print(f"\nAvailable layouts: {list(layouts2.keys())}")

# ================================================================
# Helper: get layout by name containing substring
# ================================================================
def find_layout(name_part):
    for name, idx in layouts2.items():
        if name_part in name:
            return prs2.slide_layouts[idx]
    return prs2.slide_layouts[0]

def add_cover_slide(title_text, subtitle_text, info_lines):
    """Add a title/cover slide using 标题幻灯片 layout."""
    layout = find_layout("标题")
    slide = prs2.slides.add_slide(layout)
    # Set title
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        tf.text = title_text
    # Set subtitle if placeholder exists
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # subtitle
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = subtitle_text or ""
            for line in info_lines:
                p2 = tf.add_paragraph()
                p2.text = line
    return slide

def add_content_slide(title, bullets, note=None):
    """Add a content slide using 内容页 layout."""
    layout = find_layout("内容")
    slide = prs2.slides.add_slide(layout)
    # Set title - find title placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # title
            tf = shape.text_frame
            tf.clear()
            tf.text = title
            break
    # Add bullet content in body placeholder or as new textbox
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 7:  # body
            body_shape = shape
            break
    if body_shape is None:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # subtitle/body
                body_shape = shape
                break

    if body_shape:
        tf = body_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(4)
    elif bullets:
        # No body placeholder - add text box
        left = Inches(0.92)
        top = Inches(1.85)
        width = Inches(11.5)
        height = Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.space_after = Pt(6)

    if note:
        left = Inches(0.92)
        top = Inches(6.5)
        width = Inches(11.5)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(10)

    return slide

def add_image_placeholder(slide, label, left, top, width, height):
    """Add a dashed rectangle placeholder for an image."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0x80, 0x80, 0x80)
    shape.line.width = Pt(1)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

# ================================================================
# Build Slides (following 新答辩PPT大纲.md structure)
# ================================================================

# --- Slide 1: Cover ---
add_cover_slide(
    "基于MCU的智能LED图示条幅的设计",
    "Design of an Intelligent LED Graphic Banner Based on MCU",
    [
        "聂俊宇  19220309",
        "电子科学与工程学院  电子信息工程",
        "指导教师：李秀英 副教授",
    ]
)

# --- Slide 2: 目录 ---
add_content_slide("汇报提纲", [
    "1. 选题背景与设计目标",
    "2. 设计思路和总体方案",
    "3. 硬件组成",
    "4. 软件设计 — LED显示驱动 | AI端接口 | 通信协议 | 智能绘图脚本",
    "5. 功能验证",
    "6. 总结与展望",
])

# --- Slide 3: 选题背景和意义 ---
add_content_slide("选题背景和意义", [
    "选题背景：LED图示条幅广泛应用，但传统产品显示内容固定、修改不便",
    "        人工智能技术快速发展，为LED显示智能化提供新思路",
    "",
    "设计目标：设计模块化的WS2812B LED矩阵，低功耗、易扩展",
    "         开发高效的LED驱动方案，PWM+DMA实现30fps流畅显示",
    "         构建蓝牙无线通信链路，实现AI端到LED端的数据传输",
    "         集成小智AI与MCP绘图工具，实现语音交互+智能绘图",
    "",
    "意义：提高LED图示条幅的智能化程度，探索AI+嵌入式的新交互模式",
])

# --- Slide 4: 设计思路和总体方案 ---
add_content_slide("设计思路和总体方案", [
    "系统分为两大组成部分：LED端 和 AI端",
    "",
    "LED端：",
    "  • 采用WS2812B LED组成8×8基本模块，灵活调整画幅",
    "  • AI8051U主控，复用2路PWM+DMA，分行扫描驱动",
    "  • 74HC595控制各行PMOS开关，降低静态功耗",
    "",
    "AI端：",
    "  • ESP32-S3移植小智AI，实现语音唤醒和命令识别",
    "  • 通过HC-05蓝牙模块与LED端进行无线数据传输",
    "  • 接入MCP绘图工具，供LLM调用实现智能图像生成",
])

# --- Slide 5: 系统硬件方案 ---
add_content_slide("系统硬件方案", [
    "1. LED矩阵：",
    "   WS2812B SMD 5050封装，内部集成恒流IC和RGB LED",
    "   16×16矩阵，每行独立电源线，通过PMOS开关控制",
    "",
    "2. LED控制电路：",
    "   AI8051U主控（32位，40MHz，64KB Flash + 34KB RAM）",
    "   2路PWM+DMA输出，3线SPI控制74HC595行选",
    "",
    "3. 蓝牙模块：HC-05 ×2，主从配对，UART 460800bps",
    "",
    "4. AI开发板：立创实战派ESP32-S3，集成麦克风、喇叭、触摸屏",
])

# --- Slide 6: 硬件组成 ---
add_content_slide("硬件组成", [
    "LED矩阵基本模块（8×8）：",
    "  • 8行×8列WS2812B，每行独立DI/DO和VCC/GND",
    "  • SMD封装紧凑布局，预留退耦电容焊盘",
    "  • 多模块可拼接为8×16、16×16等不同画幅",
    "",
    "LED控制电路：",
    "  • 2路PWM(P1.0/P1.2) — 驱动奇数/偶数行LED信号",
    "  • 3线595(P0.0-P0.2) — 控制16路PMOS开关",
    "  • 控制板+开关板分离，便于开发调试",
    "",
    "蓝牙模块：CSR BC417芯片 + Flash + 底板电路（3.3V稳压）",
    "AI开发板：ESP32-S3-WROOM-1-N16R8 + 音频ADC/DAC + LCD",
])

# --- Slide 7: LED显示驱动方案 ---
add_content_slide("LED显示驱动方案（软件）", [
    "WS2812B驱动层：",
    "  • PWM+DMA自动生成归零码波形，无需CPU干预",
    "  • 256×8查找表将RGB像素实时转换为PWM占空比",
    "  • 2路PWM对应奇数/偶数行，双通道并行输出",
    "",
    "图像渲染层：",
    "  • 统一渲染管线：内容类型→坐标映射→效果处理→亮度",
    "  • 多种效果：纯色、渐变、呼吸、滚动、淡入、颜色循环等",
    "  • 时间线驱动的动画调度器",
    "",
    "通信层：UART2+DMA自动接收，三态状态机解析协议包",
    "任务调度：协作式1ms调度器，管理显示刷新、动画、通信任务",
])

# --- Slide 8: AI端接口方案 ---
add_content_slide("AI端接口调度方案（软件）", [
    "双控制路径设计：",
    "",
    "路径A — 本地控制（触摸/语音关键词→SetAction直发）：",
    "  • 延迟<20ms，无需网络往返",
    "  • 适用：图案切换、时钟显示、颜色/效果切换",
    "",
    "路径B — AI绘图（自由语音→MCP工具→WebSocket回传）：",
    "  • LLM调用draw_python/show_text等工具生成图像",
    "  • 通过Debug WebSocket回传到ESP32",
    "  • 先在LCD预览，确认后通过蓝牙上传到LED端",
    "",
    "关键特性：按需启动调试服务（节省SRAM）、动画缓冲(32帧)、ACK轮询",
])

# --- Slide 9: 蓝牙通信协议 ---
add_content_slide("蓝牙通信协议设计（软件）", [
    "协议目标：可定界、完整性校验、可扩展、低带宽适配",
    "",
    "V3紧凑格式（6字节包头）：",
    "  Magic(0x47) | Flags | Seq | Cmd | PayloadLen | HeaderCRC8 | Payload | PacketCRC16",
    "",
    "三种图像格式：",
    "  • RGB332全帧: 256B — 任意全色图案",
    "  • 紧凑位图: 38B — 单色图案（32B位图+3B颜色+头部）",
    "  • 分层位图: 36B/层 — 多色叠加，≤4层可单包直传(144B)",
    "",
    "双CRC校验：header_crc8保障包头安全 + packet_crc16保障整包完整性",
    "请求-回复模型：Reply通过IS_REPLY标志+序列号+命令回显匹配原请求",
])

# --- Slide 10: 智能绘图脚本 ---
add_content_slide("智能绘图脚本与MCP集成（软件）", [
    "MCP桥接服务架构：LLM ↔ MCP(WebSocket) ↔ 桥接服务 ↔ ESP32(WebSocket) ↔ LED",
    "",
    "绘图工具集(self.screen.matrix_16x16.*)：",
    "  • draw_python — 受限Pillow绘图（AST白名单安全沙箱）",
    "  • render_prompt — 自然语言描述→图案模板",
    "  • show_text — 文字→字形→16×16位图帧序列",
    "  • show_scroll_subtitle — 长文本→离屏位图→滚动分块",
    "  • show_effect — 原生效果直连（无需逐帧渲染）",
    "  • draw_animation — 多帧动画+自动重采样（≤32帧）",
    "",
    "传输：Debug WebSocket优先（实时JSON），HTTP预览回退（PNG）",
])

# --- Slide 11: 软件方案总结/对比 ---
add_content_slide("软件方案总结", [
    "┌─────────────┬──────────────────┬──────────────────────────┐",
    "│   模块      │     核心技术     │        关键参数          │",
    "├─────────────┼──────────────────┼──────────────────────────┤",
    "│ LED显示驱动 │ PWM+DMA自动波形  │ 691kHz, 8行对扫描, 30fps │",
    "│             │ 双缓冲防撕裂     │ 32帧×36B动画缓冲         │",
    "├─────────────┼──────────────────┼──────────────────────────┤",
    "│ AI端接口    │ 双路径控制       │ 路径A<20ms, 路径B 0.5-3s │",
    "│             │ Debug WS+HTTP    │ 按需启动，节省SRAM        │",
    "├─────────────┼──────────────────┼──────────────────────────┤",
    "│ 蓝牙协议    │ V3紧凑二进制     │ 6B包头, 双CRC, ≤4层单包  │",
    "│             │ 分层位图压缩     │ 最高5×压缩比              │",
    "├─────────────┼──────────────────┼──────────────────────────┤",
    "│ 绘图脚本    │ MCP+LLM集成      │ 6种绘图工具, AST安全沙箱 │",
    "│             │ WebSocket传输    │ 32帧动画, 自动重采样      │",
    "└─────────────┴──────────────────┴──────────────────────────┘",
])

# --- Slide 12: 功能验证 ---
add_content_slide("功能验证", [
    "硬件验证：",
    "  • LED矩阵供电正常，各路PMOS开关正确",
    "  • 蓝牙配对成功，UART通信稳定（460800bps）",
    "  • 功耗测试：16×16全白时实测电流，扫描切换降低静态功耗",
    "",
    "软件验证：",
    "  • 显示效果：纯色/渐变/图案/滚动字幕/动画均正常",
    "  • 协议通信：CRC校验通过，ACK匹配正确，丢包容忍",
    "  • MCP绘图：draw_python/show_text/show_effect端到端通",
    "  • 语音控制：唤醒词检测准确，关键词→对应动作，自由绘图正常",
    "",
    "系统联调：语音「显示红心」→ AI端识别 → 绘图脚本生成 → 蓝牙传输 → LED正确显示",
])

# --- Slide 13: 总结与展望 ---
add_content_slide("总结与展望", [
    "总结：",
    "  • 完成了模块化WS2812B LED矩阵的硬件设计与制作",
    "  • 实现了PWM+DMA的高效LED驱动方案（30fps稳定显示）",
    "  • 设计了紧凑可靠的蓝牙通信协议（V3格式，双CRC校验）",
    "  • 开发了MCP智能绘图系统（LLM通过自然语言控制LED显示）",
    "  • 打通了语音→AI→绘图→蓝牙→LED的完整端到端链路",
    "",
    "创新点：",
    "  ① 双CRC轻量级协议确保蓝牙传输可靠性",
    "  ② MCP集成实现LLM控制LED绘图（标准化协议）",
    "  ③ 分层位图格式实现最高5倍传输压缩",
    "  ④ PWM+DMA自动波形生成（CPU零开销）",
    "  ⑤ 统一显示配置架构解耦本地/远程渲染",
    "",
    "展望：更高分辨率、WiFi直连、真彩色分层、手机APP控制",
])

# ================================================================
# Save as NEW file
# ================================================================
# Delete original template slides (keep only newly created ones)
# Must delete from end to avoid index shifting
orig_count = len(prs2.slides) - 13  # 13 = newly added
for i in range(orig_count - 1, -1, -1):
    rId = prs2.slides._sldIdLst[i].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs2.part.drop_rel(rId)
    prs2.slides._sldIdLst.remove(prs2.slides._sldIdLst[i])

prs2.save(OUTPUT)
print(f"\nPPT saved to: {OUTPUT}")
print(f"Total slides: {len(prs2.slides)} (original template slides removed)")

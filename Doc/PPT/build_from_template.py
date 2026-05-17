"""
Build graduation defense PPT from 本科毕业论文答辩模板.pptx template.
Follows thesis chapter structure. Creates NEW file.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

TEMPLATE = r"D:\GraduationProject\Doc\PPT\本科毕业论文答辩模板.pptx"
OUTPUT = r"D:\GraduationProject\Doc\PPT\毕业答辩-聂俊宇-模板版.pptx"

prs = Presentation(TEMPLATE)
W = prs.slide_width
H = prs.slide_height

print(f"Slide size: {W/914400:.1f}x{H/914400:.1f} in, {len(prs.slides)} slides")

# ================================================================
# Helpers
# ================================================================
def fill_placeholder(slide, ph_idx, text, font_size=None):
    """Fill a placeholder with text."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            return tf

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return tf

def add_bullets(slide, left, top, width, height, items, font_size=13):
    """Add bulleted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.space_after = Pt(4)
    return tf

def add_img_placeholder(slide, label, left, top, width, height):
    """Dashed rectangle with label for image placeholder."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)
    shape.line.width = Pt(1.2)
    shape.line.dash_style = 2
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)

def set_title(slide, title_text):
    """Set the Title Only placeholder (idx 0)."""
    return fill_placeholder(slide, 0, title_text, 28)

def set_section(slide, main_text, sub_text=""):
    """Set Section Header placeholders."""
    fill_placeholder(slide, 0, main_text, 30)
    if sub_text:
        fill_placeholder(slide, 1, sub_text, 16)

# ================================================================
# Edit each slide
# ================================================================
s = prs.slides

# --- Slide 1: Title Slide (Cover) ---
print("Slide 1: Cover...")
sl = s[0]
# Match placeholders by position (all have same type=14, different positions)
for shape in sl.placeholders:
    y_in = shape.top / 914400
    if y_in < 2.0:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "基于MCU的智能LED图示条幅的设计"
    elif y_in < 4.5:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "Design of an Intelligent LED Graphic Banner Based on MCU"
    elif shape.left / 914400 < 6.0:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "答辩人：聂俊宇"
    else:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "指导教师：李秀英 副教授"

# --- Slide 2: Blank → Table of Contents ---
print("Slide 2: TOC...")
sl = s[1]
toc_items = [
    "1. 选题背景与意义",
    "2. 总体设计方案",
    "3. 硬件组成",
    "4. 软件设计",
    "5. 功能验证",
    "6. 总结与展望",
]
add_textbox(sl, 1.5, 0.8, 10, 0.6, "汇报提纲", 28, True)
add_bullets(sl, 2.5, 1.8, 8, 4.5, toc_items, 18)

# --- Slide 3: Section Header — 第1章 绪论 ---
print("Slide 3: Section...")
set_section(s[2], "第1章  绪论", "研究背景、国内外现状与研究内容")

# --- Slide 4: Title Only — 研究背景 ---
print("Slide 4: 研究背景...")
set_title(s[3], "研究背景")
add_bullets(s[3], 0.7, 1.3, 7.5, 4.5, [
    "LED图示条幅广泛应用于节日庆典、广告宣传、室内装饰等场景",
    "传统LED图示条幅产品显示内容固定，修改不便，缺乏灵活性",
    "2022年底ChatGPT发布以来，AI技术日新月异，各行业与AI深度融合",
    "2025年国务院印发\"人工智能+\"行动意见，AI上升为国家战略",
    "\"AI+\"浪潮中，LED显示正向\"双向交互\"与\"智能化\"转型",
    "小智AI开源项目为嵌入式设备提供语音交互+工具扩展能力",
], 14)
add_img_placeholder(s[3], "传统LED条幅照片", 8.7, 1.3, 3.8, 2.2)
add_img_placeholder(s[3], "AI+LED概念图", 8.7, 3.8, 3.8, 2.0)

# --- Slide 5: Title Only — 研究意义 ---
print("Slide 5: 研究意义...")
set_title(s[4], "研究意义")
add_bullets(s[4], 0.7, 1.3, 11.5, 5.2, [
    "理论意义：",
    "  • 探索AI大模型与嵌入式LED显示系统的融合方法",
    "  • 提出紧凑二进制协议设计，为资源受限设备通信提供参考",
    "  • 验证PWM+DMA方案的WS2812驱动可行性",
    "",
    "应用价值：",
    "  • 会议场景：语音快速修改会议主题和时间，提高信息传达效率",
    "  • 教育场景：作为单词学习工具，显示英文单词和图像",
    "  • 家居场景：多功能时钟，语音设置闹钟、计时等",
    "  • 商用场景：智能导购指示牌、沉浸式互动体验装置",
], 13)

# --- Slide 6: Title Only — 研究内容 ---
print("Slide 6: 研究内容...")
set_title(s[5], "本文研究内容")
add_bullets(s[5], 0.7, 1.3, 7.0, 5.2, [
    "1. 模块化WS2812B LED矩阵的设计与制作",
    "   • 8×8基本模块，16×16画幅，行扫描架构",
    "",
    "2. 高效LED显示驱动方案开发",
    "   • AI8051U + PWM+DMA，30fps流畅显示",
    "   • 协作式任务调度器，多任务管理",
    "",
    "3. 蓝牙通信协议设计",
    "   • 轻量二进制协议，双CRC校验，3种图像格式",
    "",
    "4. AI端接口与智能绘图系统开发",
    "   • 双控制路径（本地+远程），MCP工具集成",
    "   • 语音交互控制LED显示和智能绘图",
], 13)
add_img_placeholder(s[5], "系统整体效果图", 8.3, 1.3, 4.2, 5.2)

# --- Slide 7: Section Header — 第2章 总体方案 ---
print("Slide 7: Section...")
set_section(s[6], "第2章  总体方案", "设计目标、思路与软硬件方案")

# --- Slide 8: Title Only — 设计目标与思路 ---
print("Slide 8: 设计目标与思路...")
set_title(s[7], "设计目标与总体思路")
add_bullets(s[7], 0.7, 1.3, 11.5, 5.2, [
    "设计目标：",
    "  硬件 — 模块化WS2812B LED矩阵 + AI8051U控制电路，低功耗、易扩展",
    "  软件 — 流畅LED驱动(30fps) + 蓝牙传输 + 语音AI + 智能绘图",
    "",
    "总体思路（两大组成部分）：",
    "  LED端 — WS2812B矩阵 + AI8051U控制电路 + 行扫描驱动",
    "         复用2路PWM+DMA + 74HC595行选 + UART协议接收",
    "  AI端  — ESP32-S3 + 小智AI + HC-05蓝牙 + MCP绘图工具",
    "         语音交互 → LED接口 → 蓝牙发包 + WebSocket绘图接收",
], 13)

# --- Slide 9: Title Only — 硬件方案与选型 ---
print("Slide 9: 硬件方案...")
set_title(s[8], "硬件方案与器件选型")
add_bullets(s[8], 0.7, 1.3, 7.2, 5.2, [
    "硬件方案：",
    "  • LED矩阵: WS2812B × 256 (16×16), 8行对扫描",
    "  • LED控制: AI8051U + 2路PWM+DMA + 74HC595行选",
    "  • 蓝牙通信: HC-05 ×2 主从配对, UART 460800bps",
    "  • AI开发板: 立创实战派ESP32-S3, 集成音频+触摸",
    "",
    "器件选型理由：",
    "  • WS2812B: SMD5050, 集成IC+RGB, 单线级联, 外围简单",
    "  • AI8051U: 32位, 40MHz, 64KB Flash+34KB RAM",
    "    成本1.9元 vs STM32/GD32 2.2元以上",
    "  • HC-05: 经典蓝牙SPP, 支持AT指令, 最高1382400bps",
    "  • ESP32-S3: 240MHz, AI向量指令, Wi-Fi+BLE",
], 13)
add_img_placeholder(s[8], "硬件选型对比照片\n(WS2812B+AI8051U+HC-05+ESP32)", 8.3, 1.3, 4.2, 5.2)

# --- Slide 10: Title Only — 软件方案 ---
print("Slide 10: 软件方案...")
set_title(s[9], "软件方案概述")
add_bullets(s[9], 0.7, 1.3, 11.5, 5.2, [
    "LED显示驱动方案：Timer1 1ms行扫描 + PWM+DMA自动波形 + 协作式任务调度(30fps)",
    "  WS2812B驱动层 → 帧缓存编码层 → 图像渲染层 → 动作控制层 → 通信接入层",
    "",
    "AI端接口方案：双控制路径",
    "  路径A: 本地触摸/语音关键词 → SetAction直发蓝牙 (<20ms)",
    "  路径B: 自由语音 → MCP工具调用 → 主机绘图 → WebSocket回传 → 蓝牙上传",
    "",
    "蓝牙通信协议：V3紧凑格式(6B包头), 双CRC校验, 三种图像格式",
    "  RGB332全帧(256B) / 紧凑位图(38B) / 分层位图(36B/层)",
    "",
    "智能绘图脚本：MCP桥接 + 6种绘图工具 + AST白名单安全沙箱",
    "  draw_python / render_prompt / show_text / show_scroll_subtitle / show_effect / draw_animation",
], 12)

# --- Slide 11: Section Header — 第3-4章 硬件与软件 ---
print("Slide 11: Section...")
set_section(s[10], "第3-4章  硬件组成与软件设计", "系统各模块的硬件结构与软件实现逻辑")

# --- Slide 12: Title Only — 硬件组成 ---
print("Slide 12: 硬件组成...")
set_title(s[11], "硬件组成")
add_bullets(s[11], 0.7, 1.3, 7.5, 3.5, [
    "LED矩阵 (16×16)：",
    "  • 正面: 各行独立VCC(粉) + 公共GND(黄) + DI/DO信号(红)",
    "  • 背面: 公共GND(黄) + 外部信号连线(蓝) + 终端接地电阻",
    "  • 每灯珠预留退耦电容焊盘，保证电源稳定性",
    "",
    "LED控制电路：",
    "  • 控制板: AI8051U接口 + 74HC595 + 外部接线端子",
    "  • 开关板: 16路PMOS高侧开关, 逐行控制电源降低静态功耗",
    "  • 静态电流优化: 64LED全熄实测43mA → 行关闭后显著降低",
    "",
    "蓝牙模块: CSR BC417 + Flash + 底板电路(3.3V稳压)",
    "AI开发板: ESP32-S3-WROOM-1 + 音频ADC/DAC + LCD触摸屏",
], 13)
add_img_placeholder(s[11], "LED矩阵PCB照片(正反面)", 8.7, 1.3, 3.8, 2.2)
add_img_placeholder(s[11], "控制电路原理图/PCB", 8.7, 3.8, 3.8, 2.2)

# --- Slide 13: Title Only — LED显示驱动 ---
print("Slide 13: LED显示驱动...")
set_title(s[12], "LED显示驱动方案（软件核心）")
add_bullets(s[12], 0.7, 1.3, 6.8, 5.2, [
    "WS2812B驱动层（PWM+DMA）：",
    "  • PWM频率~691kHz (33.18MHz/48), 1.447us/bit",
    "  • 256×8查找表: RGB像素→PWM占空比 实时编码",
    "  • 双通道(CH1偶行/CH2奇行) 双行交织DMA缓冲",
    "  • 8步完成一帧(8行对×~1ms), 约30fps刷新",
    "",
    "图像渲染层(draw_drv.c)：",
    "  • 统一管线: 内容类型→坐标映射→颜色重映射→效果→亮度",
    "  • 效果集合: 纯色/渐变/呼吸/滚动/淡入淡出/颜色循环/行显隐",
    "  • 统一时间线驱动的动画调度器",
    "",
    "通信层: UART2+DMA 192B环形缓冲 + 三态收包状态机",
    "任务调度: 1ms协作式调度器(≤8任务, 无抢占, 低RAM)",
], 12)
add_img_placeholder(s[12], "PWM编码管线流程图", 7.9, 1.3, 4.6, 2.5)
add_img_placeholder(s[12], "示波器PWM波形截图\n(0码/1码/RESET)", 7.9, 4.1, 4.6, 2.5)

# --- Slide 14: Title Only — AI端接口+协议 ---
print("Slide 14: AI接口+协议...")
set_title(s[13], "AI端接口与蓝牙通信协议")
add_bullets(s[13], 0.7, 1.3, 6.8, 5.2, [
    "AI端接口调度 (ESP32-S3 + 小智AI)：",
    "  • 路径A(直发): 触摸UI/语音关键词 → SetAction → 蓝牙 (<20ms)",
    "  • 路径B(绘图): 语音 → LLM调用MCP工具 → WebSocket回传 → 蓝牙上传",
    "  • 按需启动Debug服务: 节省SRAM, listening阶段不常驻",
    "",
    "蓝牙通信协议 (V3紧凑格式)：",
    "  • 6B包头: Magic(0x47)+Flags+Seq+Cmd+Len+CRC8",
    "  • 双CRC: header_crc8(包头安全) + packet_crc16(整包完整)",
    "  • 3种图像格式: RGB332(256B) / 紧凑位图(38B) / 分层位图(36B/层)",
    "  • ≤4层144B单包直传 (效率84% vs 分块59%)",
    "  • 4类命令: 参数控制/整帧/动画/轻量图像, 20+条命令",
], 12)
add_img_placeholder(s[13], "双路径控制流图", 8.0, 1.3, 4.5, 2.3)
add_img_placeholder(s[13], "协议包字节布局图", 8.0, 3.9, 4.5, 2.0)

# --- Slide 15: Title Only — 绘图脚本 ---
print("Slide 15: 绘图脚本...")
set_title(s[14], "智能绘图脚本与MCP集成")
add_bullets(s[14], 0.7, 1.3, 6.8, 5.2, [
    "MCP桥接服务架构：",
    "  LLM ↔ MCP(WebSocket) ↔ 桥接服务 ↔ ESP32(Debug WS) ↔ LED",
    "",
    "绘图工具集 (self.screen.matrix_16x16.*)：",
    "  • draw_python: AST白名单安全绘图 (受限Pillow执行)",
    "  • render_prompt: 自然语言 → 图案模板",
    "  • show_text: 文字 → 字形 → 16×16位图帧序列",
    "  • show_scroll_subtitle: 长文本 → 离屏位图 → 滚动分块",
    "  • show_effect: 原生效果直连 (无需逐帧渲染)",
    "  • draw_animation: 多帧动画 + 自动重采样 (≤32帧)",
    "",
    "传输: Debug WebSocket优先(实时JSON), HTTP预览回退(PNG)",
    "安全: AST白名单, 仅允许安全绘图函数, 禁止import/eval/文件IO",
], 12)
add_img_placeholder(s[14], "MCP架构流程图", 7.9, 1.3, 4.6, 2.8)
add_img_placeholder(s[14], "MCP工具列表截图", 7.9, 4.3, 4.6, 2.2)

# --- Slide 16: Section Header — 第5章 功能验证 ---
print("Slide 16: Section...")
set_section(s[15], "第5章  功能验证", "系统硬件测试与软件功能验证结果")

# --- Slide 17: Title Only — 测试结果 ---
print("Slide 17: 测试结果...")
set_title(s[16], "功能验证结果")
add_bullets(s[16], 0.7, 1.3, 7.5, 5.2, [
    "硬件验证：",
    "  • LED矩阵供电正常, 16路PMOS开关切换正确",
    "  • 蓝牙配对成功, UART 460800bps稳定通信",
    "  • 功耗测试: 全熄<1mA(行开关关闭), 全白实测工作电流",
    "",
    "显示效果：纯色/渐变/图案/滚动字幕/动画 30fps流畅",
    "",
    "协议通信：CRC校验通过率100%, ACK匹配正确, 丢包容忍",
    "",
    "语音控制：唤醒词检测准确, 关键词→对应动作正常",
    "",
    "MCP智能绘图：",
    "  • draw_python/show_text/show_effect 端到端通过",
    "  • AST安全沙箱有效, 危险操作被拦截",
    "  • 语音\"显示红心\"→ LED正确显示红色爱心",
], 13)
add_img_placeholder(s[16], "系统工作实物照片\n(多种显示效果)", 8.7, 1.3, 3.8, 2.5)
add_img_placeholder(s[16], "测试结果截图\n(串口日志/MCP工具)", 8.7, 4.1, 3.8, 2.5)

# --- Slide 18: Closing — 致谢 ---
print("Slide 18: Closing...")
sl = s[17]
for shape in sl.placeholders:
    y_in = shape.top / 914400
    if y_in < 3.0:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "感谢聆听  请老师批评指正"
    elif shape.left / 914400 < 6.0:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "答辩人：聂俊宇"
    else:
        shape.text_frame.clear()
        shape.text_frame.paragraphs[0].text = "电子科学与工程学院"

# ================================================================
# Save
# ================================================================
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
